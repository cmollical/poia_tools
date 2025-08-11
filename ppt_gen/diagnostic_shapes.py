#!/usr/bin/env python
"""Diagnostic script – enumerate placeholder shapes

This utility duplicates the first slide of the Roadmap_Preview_Deep_Dive template
and replaces every text-box shape with its own index and shape name so you can
visually identify which shapes to target in automation scripts.
"""
from __future__ import annotations
import os
from pptx import Presentation

# ---------------------------------------------------------------------------
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PPT = os.path.join(
    SCRIPT_DIR,
    "ppt_templates",
    "Roadmap_Preview_Deep_Dive.pptx",
)
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "generated")
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_PPT = os.path.join(OUTPUT_DIR, "Diagnostic_Placeholders.pptx")


def main() -> None:
    if not os.path.exists(TEMPLATE_PPT):
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PPT}")

    prs = Presentation(TEMPLATE_PPT)
    slide = prs.slides[0]  # operate on first slide directly

    # Replace text with diagnostics
    for idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = f"{idx} | {shape.name}"

    prs.save(OUTPUT_PPT)
    print(f"Diagnostic presentation saved to: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
