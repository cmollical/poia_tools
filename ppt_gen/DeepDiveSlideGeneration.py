#!/usr/bin/env python
"""
Targeted fix for both timeframe and title/subtitle in Roadmap Preview slides.
"""

import os
import sys
import logging
import traceback
import re
from datetime import datetime
from copy import deepcopy
import pandas as pd
from openai import AzureOpenAI
from pptx import Presentation
from pptx.enum.text import PP_ALIGN  # Import for text alignment

# Use the new dbUtils module that matches Ask Amy's approach
from dbUtils import execute_snowflake_query

# Database configuration
DATABASE = "CORPANALYTICS_BUSINESS_PROD"
SCHEMA = "SCRATCHPAD_PRDPF"

# AthenaGPT configuration
AGPT_API_VERSION = "2023-12-01-preview"
AGPT_ENDPOINT = "https://athena-gpt-prod-e.openai.azure.com/"
AGPT_MODEL = "gpt-4o-mini"

def _get_release_index(release_code: str) -> tuple[int, int, int, int]:
    """
    Converts a release code (e.g., '24.07') into structured data, including a
    sequential quarterly index for comparisons.

    Returns:
        A tuple containing (year_short, month, quarter, overall_index).
    """
    try:
        year_short, month = map(int, release_code.split('.'))
        
        if not (1 <= month <= 12):
            raise ValueError("Month must be between 1 and 12")

        if 1 <= month <= 3:
            quarter = 1
        elif 4 <= month <= 6:
            quarter = 2
        elif 7 <= month <= 9:
            quarter = 3
        else:  # 10-12
            quarter = 4
            
        # A baseline year (e.g., 2020) is used to create a consistent sequential index.
        overall_index = (year_short - 20) * 4 + (quarter - 1)
        
        return year_short, month, quarter, overall_index
    except (ValueError, IndexError) as e:
        logger.error(f"Invalid release code format: '{release_code}'. {e}")
        raise ValueError(f"Invalid release code format: '{release_code}'") from e


def _format_timeframe(release_code: str) -> str:
    """
    Formats a release code into a human-readable timeframe based on its
    difference from the current release period.

    - "Recently Released": For releases in a past quarter.
    - "Season YYYY": For the current or next quarterly release.
    - "Season YYYY or later": For releases >1 quarter away in the same year.
    - "Future Release": For releases in a future year.
    """
    if not release_code or pd.isna(release_code):
        return "TBD"

    release_code = str(release_code).strip().replace('*', '')
    print(f"DEBUG: Formatting timeframe for release code: '{release_code}'")

    try:
        feature_year_short, feature_month, _, feature_index = _get_release_index(release_code)
        
        now = datetime.now()
        current_release_code = f"{now.year - 2000}.{now.month:02d}"
        current_year_short, _, _, current_index = _get_release_index(current_release_code)

        index_diff = feature_index - current_index
        
        year = 2000 + feature_year_short
        if 1 <= feature_month <= 3: season = "Winter"
        elif 4 <= feature_month <= 6: season = "Spring"
        elif 7 <= feature_month <= 9: season = "Summer"
        else: season = "Fall"
        base_format = f"{season} {year}"

        if index_diff < 0:
            return "Recently Released"
        elif index_diff == 0:
            return base_format
        else:  # index_diff > 0
            if feature_year_short > current_year_short:
                return "Future Release"
            
            if index_diff > 1:
                return f"{base_format} or later"
            else:  # index_diff == 1
                return base_format

    except (ValueError, IndexError) as e:
        logger.warning(f"Could not parse release code '{release_code}'. Defaulting to 'TBD'. Error: {e}")
        return "TBD"
    
    return "TBD"


# AthenaGPT integration constants
AGPT_MODEL: str = "gpt-4o-mini-2024-07-18"
AGPT_ENDPOINT: str = os.getenv("AGPT_ENDPOINT", "https://athenagpt-uat.tools.athenahealth.com/api/public/oai")
AGPT_API_VERSION: str = os.getenv("AGPT_API_VERSION", "2025-01-01-preview")


def _athenagpt_complete(prompt: str, temperature: float = 0.7, max_tokens: int = 400) -> str:
    """Call AthenaGPT using AzureOpenAI SDK and return the assistant response text."""
    # Debug environment variables
    logger.info("Checking for AthenaGPT API key in environment variables")
    api_key = os.getenv("AGPT_API") or os.getenv("ATHENAGPT_API_KEY") or os.getenv("AGPT_KEY") or os.getenv("AZURE_OPENAI_API_KEY")
    
    # Print available environment variables for debugging (without exposing sensitive values)
    env_vars = [k for k in os.environ.keys() if 'API' in k.upper() or 'KEY' in k.upper() or 'AGPT' in k.upper() or 'GPT' in k.upper()]
    logger.info(f"Available environment variables that might contain API keys: {env_vars}")
    
    if not api_key:
        logger.error("AthenaGPT API key not found in any of the expected environment variables")
        raise RuntimeError("AthenaGPT API key not found – set AGPT_API / ATHENAGPT_API_KEY / AGPT_KEY")
    try:
        client = AzureOpenAI(
            api_version=AGPT_API_VERSION,
            api_key=api_key,
            azure_endpoint=AGPT_ENDPOINT,
        )
        completion = client.chat.completions.create(
            model=AGPT_MODEL,
            messages=[
                {"role": "system", "content": "You are a senior product manager with expertise in healthcare technology and athenahealth products."},
                {"role": "user", "content": prompt},
            ],
            temperature=temperature,
            max_tokens=max_tokens,
        )
        return completion.choices[0].message.content.strip()
    except Exception as exc:
        raise RuntimeError(f"athenaGPT API error: {exc}")

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Script directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Snowflake configuration
SNOWFLAKE_CFG = dict(
    account   = "athenahealth",
    user      = "SVC_JIR_PROPS",
    password  = os.getenv("SERVICE_PASS", ""),
    role      = "CORPANALYTICS_BDB_PRDPF_PROD_RW",
    warehouse = "CORPANALYTICS_BDB_PRDPF_WH_READWRITE_PROD",
    database  = "CORPANALYTICS_BUSINESS_PROD",
    schema    = "SCRATCHPAD_PRDPF",
    autocommit=True,
)

def parse_feature_keys(feature_keys_input: str) -> list:
    """Parse user input to extract feature keys."""
    if not feature_keys_input:
        return []
    keys = re.split(r'[\s,;\n]+', feature_keys_input)
    return [key.strip().upper() for key in keys if key.strip() and key.upper().startswith("FEATURE-")]

def parse_release_codes(release_codes_input: str) -> list:
    """
    Parse user input to extract release codes (e.g., '26.03').
    Handles codes with or without surrounding asterisks.
    """
    if not release_codes_input:
        return []
    
    potential_codes = re.split(r'[\s,;\n]+', release_codes_input)
    valid_codes = []
    
    for code in potential_codes:
        # Strip whitespace and any surrounding asterisks before validation
        sanitized_code = code.strip().strip('*')
        if re.match(r'^\d{2}\.\d{2}$', sanitized_code):
            valid_codes.append(sanitized_code)
            
    return valid_codes

def load_feature_keys_by_release(release_codes: list) -> list:
    """
    Load all feature keys associated with the given release codes, filtered for
    features designated for deep dive slide creation.
    """
    if not release_codes:
        return []
    
    logger.info(f"Loading feature keys for releases: {release_codes}")
    
    # Sanitize release codes and build a series of LIKE conditions for the SQL query.
    sanitized_codes = [code.strip().strip('*') for code in release_codes]
    release_like_conditions = " OR ".join([f"TARGET_GA_RELEASE LIKE '%{code}%'" for code in sanitized_codes])

    # Define the specific artifact value to filter by
    artifact_filter = "External Roadmap: Create deep dive slide"
    
    query = f"""
    SELECT DISTINCT FEATURE_KEY
    FROM CORPANALYTICS_BUSINESS_PROD.SCRATCHPAD_PRDPF.FEATURE_API_FULL
    WHERE 
        ({release_like_conditions})
        AND INCLUDE_IN_ROADMAP_ARTIFACTS LIKE '%{artifact_filter}%'
    """

    # Enhanced diagnostic logging
    logger.critical("="*80)
    logger.critical("EXECUTING THE FOLLOWING SQL QUERY FOR RELEASE-BASED SLIDE GENERATION:")
    logger.critical(query)
    logger.critical("="*80)

    try:
        results = execute_snowflake_query(query)
        df = pd.DataFrame(results)
        if df.empty:
            logger.warning(f"No features marked for deep dive slides found for releases: {release_codes}")
            return []
        feature_keys = df['FEATURE_KEY'].tolist()
        logger.info(f"Found {len(feature_keys)} features for deep dive slides for the given releases.")
        return feature_keys
    except Exception as e:
        logger.error(f"Error loading feature keys by release: {e}")
        logger.error(traceback.format_exc())
        return []

def format_feature_keys_for_sql(feature_keys: list) -> str:
    """Format feature keys for SQL IN clause."""
    formatted_keys = ', '.join([f"'{key}'" for key in feature_keys])
    return formatted_keys

def load_feature_data(feature_keys: list) -> pd.DataFrame:
    """Load feature data from Snowflake for the given feature keys."""
    if not feature_keys:
        logger.error("No feature keys provided")
        return pd.DataFrame()  # Empty DataFrame
    
    logger.info(f"Loading data for {len(feature_keys)} features")
    
    # SQL query for fetching feature data
    SQL_FEATURE_BY_KEY = """
    SELECT 
        FEATURE_KEY,
        CLIENT_FACING_FEATURE_NAME,
        EXTERNALROADMAPTIMEFRAME,
        EXTERNALROADMAPLANGUAGE,
        WHAT_VALUE_DOES_IT_DELIVER,
        WHAT_IS_YOUR_FEATURE,
        TARGET_GA_RELEASE,
        INDEX_RELEASES_AWAY,
        INCLUDE_IN_ROADMAP_ARTIFACTS,
        PO,
        PO_EMAIL
    FROM CORPANALYTICS_BUSINESS_PROD.SCRATCHPAD_PRDPF.FEATURE_API_FULL
    WHERE FEATURE_KEY IN ({feature_keys})
    """
    
    # Prepare SQL
    formatted_keys = format_feature_keys_for_sql(feature_keys)
    base_query = SQL_FEATURE_BY_KEY.format(feature_keys=formatted_keys)
    
    try:
        # Use dbUtils to execute the query
        logger.info("Running query for basic feature info")
        results = execute_snowflake_query(base_query)
        df = pd.DataFrame(results)
        df['ATHENAGPT_BULLETS'] = None # Initialize column for generated bullets
        logger.info(f"Query returned {len(df)} rows")
        
        # If no features found, return empty DataFrame
        if len(df) == 0:
            logger.warning(f"No features found for keys: {feature_keys}")
            return pd.DataFrame()
        
        # Generate bullet points using AthenaGPT instead of Snowflake Cortex
        for idx, feature in df.iterrows():
            feature_key = feature['FEATURE_KEY']
            
            # Build prompt for AthenaGPT
            prompt = (
                f"FEATURE_NAME: {feature['CLIENT_FACING_FEATURE_NAME']}\n"
                f"Roadmap Language: {feature['EXTERNALROADMAPLANGUAGE']}\n"
                f"Feature Value Statement: {feature['WHAT_VALUE_DOES_IT_DELIVER']}\n"
                f"Feature Description: {feature['WHAT_IS_YOUR_FEATURE']}\n\n"
                "Create 3 concise, high-impact bullet points that would be compelling to healthcare executives and providers. "
                "Each bullet should focus on a different aspect (Clinical Impact, Workflow Efficiency, Financial/Strategic).\n"
                "IMPORTANT GUIDELINES:\n"
                "- Start each bullet with a strong verb or adverb (e.g., 'Automate', 'Empower', 'Reduce').\n"
                "- Use sentence case and keep each bullet to roughly one slide line (≈ 15 words).\n"
                "- Do not include bullet symbols or numbering; output just the 3 lines separated by newlines."
            )
            try:
                gpt_response = _athenagpt_complete(prompt)
                bullet_points = parse_bullet_points(gpt_response)
                logger.info(f"AthenaGPT bullets for {feature_key}: {bullet_points}")
                df.at[idx, 'ATHENAGPT_BULLETS'] = bullet_points
            except Exception as gpt_err:
                logger.error(f"AthenaGPT failed for {feature_key}: {gpt_err}")
                df.at[idx, 'ATHENAGPT_BULLETS'] = []
        
        # Always return the DataFrame
        return df
            
    except Exception as e:
        logger.error(f"Error loading feature data: {str(e)}")
        logger.error(traceback.format_exc())
        return pd.DataFrame()  # Return empty DataFrame on error

def parse_bullet_points(description: str) -> list:
    """Parse bullet points from the Cortex response."""
    if not description or description == 'No description available':
        return []
    
    # Clean up the description - remove quotes and extra spaces
    clean_desc = description.strip().strip('"\'').strip()
    
    # Split by newlines to get individual bullet points
    bullet_points = [line.strip() for line in clean_desc.split('\n') if line.strip()]
    
    # Remove any bullet symbols or numbering
    import re
    bullet_pattern = r'^[\*\-•◦‣⁃⁌⁍⦾⦿]|\d+[\.\)]\s*'
    
    parsed_bullets = []
    for bullet in bullet_points:
        # Remove any bullet symbols or numbering
        clean_bullet = re.sub(bullet_pattern, '', bullet).strip()
        if clean_bullet:
            parsed_bullets.append(clean_bullet)
    
    return parsed_bullets

def replace_text_in_shape(shape, old_text, new_text):
    """Replace text in a shape while preserving formatting."""
    if not hasattr(shape, 'text_frame'):
        return False
    
    # Check if the shape contains the old text
    if shape.text == old_text:
        # Set the text directly
        shape.text = new_text
        return True
    
    # If we didn't find an exact match, log the text for debugging
    logger.debug(f"Shape text '{shape.text}' didn't match '{old_text}'")
    
    return False

def find_shape_by_name(slide, name):
    """Find a shape by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    logger.warning(f"Shape with name '{name}' not found on slide.")
    return None

def timeframe_title_fix(template_path, output_path, feature_keys):
    """
    Generate PowerPoint slides by creating a new slide for each feature and populating the placeholders.
    """
    logger.info(f"Starting slide generation with template: {template_path}")
    
    # Load feature data
    feature_data = load_feature_data(feature_keys)
    # Check if feature_data is None or empty
    if feature_data is None or (hasattr(feature_data, 'empty') and feature_data.empty):
        raise ValueError("No feature data found for the given keys.")

    prs = Presentation(template_path)
    template_slide = prs.slides[0]
    
    for index, feature in feature_data.iterrows():
        feature_key = feature.get('FEATURE_KEY', 'N/A')
        logger.info(f"Processing feature {index + 1}/{len(feature_data)}: {feature_key}")

        new_slide = prs.slides.add_slide(template_slide.slide_layout)

        try:
            # Get feature data
            client_name = feature.get('CLIENT_FACING_FEATURE_NAME', '')
            roadmap_lang = feature.get('EXTERNALROADMAPLANGUAGE', '')
            timeframe = _format_timeframe(feature.get('TARGET_GA_RELEASE', ''))
            bullet_points = feature.get('ATHENAGPT_BULLETS', [])
            if not isinstance(bullet_points, list): bullet_points = []
            while len(bullet_points) < 3: bullet_points.append("")

            # Parse title and subtitle
            raw_roadmap = (roadmap_lang or '').strip()
            title, subtitle = ( [part.strip() for part in raw_roadmap.split(' - ', 1)] if ' - ' in raw_roadmap else (raw_roadmap, ''))
            if not title: title = client_name or 'Unnamed Feature'

            # --- Populate Placeholders --- 
            for shape in new_slide.placeholders:
                ph_name = shape.name
                
                # Title
                if 'Title' in ph_name or ph_name == 'Title 6':
                    shape.text_frame.text = title
                    shape.text_frame.paragraphs[0].font.bold = True
                # Subtitle
                elif ph_name == 'Text Placeholder 1':
                    shape.text_frame.text = subtitle
                # Timeframe
                elif ph_name == 'Text Placeholder 7':
                    shape.text_frame.text = timeframe
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                # Bullets
                elif ph_name == 'Text Placeholder 4':
                    shape.text_frame.text = bullet_points[0]
                elif ph_name == 'Text Placeholder 3':
                    shape.text_frame.text = bullet_points[1]
                elif ph_name == 'Text Placeholder 2':
                    shape.text_frame.text = bullet_points[2]

            # --- Add Notes --- 
            notes_slide = new_slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.clear()
            p = notes_frame.add_paragraph()
            p.text = (
                f"Feature ID: {feature.get('FEATURE_KEY', 'N/A')}\n"
                f"Target GA Release: {feature.get('TARGET_GA_RELEASE', 'N/A')}\n"
                f"Product Owner: {feature.get('PO', 'N/A')}\n"
                f"PO Email: {feature.get('PO_EMAIL', 'N/A')}"
            )
            logger.info(f"Successfully wrote notes for {feature_key}")

        except Exception as e:
            logger.error(f"Error processing feature {feature_key}: {str(e)}")
            logger.error(traceback.format_exc())
            continue

    # Delete the original template slide
    if len(prs.slides) > len(feature_data):
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
        logger.info("Template slide deleted.")

    prs.save(output_path)
    logger.info(f"Finished generating presentation: {output_path}")
    return output_path

if __name__ == "__main__":
    print("Timeframe and Title Fix for Roadmap Preview Slides")
    print("=================================================")
    
    feature_keys = ["FEATURE-25612"]
    if len(sys.argv) > 1:
        feature_keys = parse_feature_keys(sys.argv[1])
    
    template_path = os.path.join(SCRIPT_DIR, "ppt_templates", "Roadmap_Preview_Deep_Dive.pptx")
    output_path = os.path.join(SCRIPT_DIR, "generated", f"Roadmap_Preview_Deep_Dive_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    
    try:
        result = timeframe_title_fix(template_path, output_path, feature_keys)
        print(f"\nSuccess! Presentation saved to: {result}")
    except Exception as e:
        print(f"\nError: {str(e)}")
        print("See log output above for details.")
