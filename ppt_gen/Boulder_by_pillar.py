#!/usr/bin/env python
"""
Build Boulder Executive Update deck.

• Template slide   : slide index 0 in BOULDER_PPT_TEMPLATE.pptx
• Table layout     : Boulder_Group | Health_Check | Boulder_Name |
                     25.3 | 25.7 | 25.11 | 26.X+ | Latest_Update
• One slide per Pillar, rows inserted/deleted to fit exactly.
• Boulder_Group cells vertically merged for contiguous rows.
• Latest_Update truncated to 500 words (+ “ … ”); now 8 pt font.
• Feature counts fixed (strips asterisks from TARGET_GA_RELEASE).
"""

from __future__ import annotations
import os, re, logging, traceback
from copy import deepcopy
from datetime import datetime, timezone

import pandas as pd

# Use the new dbUtils module that matches Ask Amy's approach
from dbUtils import execute_snowflake_query
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement    
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

# ───────── Configuration ─────────
# SCRIPT_DIR: Folder where this script lives.
# TEMPLATE  : Path to the PowerPoint template slide deck.  
#             ➡️  CHANGE this if you want to use a different .pptx as the starting point.
# LOG_DIR   : Where logs are written (ensure the user running the script has write access).
#
# Tip: If you move or rename the template file, update TEMPLATE accordingly.
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE   = os.path.join(SCRIPT_DIR, "ppt_templates", "BOULDER_PPT_TEMPLATE.pptx")
LOG_DIR    = os.path.join(SCRIPT_DIR, "Logs")
os.makedirs(LOG_DIR, exist_ok=True)
logging.basicConfig(
    filename=os.path.join(
        LOG_DIR, f"boulder_ppt_{datetime.now(timezone.utc):%Y%m%d_%H%M%S}.log"
    ),
    level=logging.INFO,
    format="%(asctime)s — %(levelname)s — %(message)s",
)
logger = logging.getLogger(__name__)

# Database configuration
DATABASE = "CORPANALYTICS_BUSINESS_PROD"
SCHEMA = "SCRATCHPAD_PRDPF"

GA_COLS  = ["25.03", "25.07", "25.11", "26.X+"]  # Default columns
# Columns to display for GA release counts in the table.  
# ➡️  Edit GA_COLS to add/remove releases or change their order.
MAX_WORDS = 500  # Max words to keep for the Latest_Update column

def truncate_words(text: str | None, limit: int = 500) -> str:
    """Clean markup, truncate to *limit* words, and return ellipsis if truncated."""
    if not text:
        return ""
    # Remove XML/Excel markup like _x000D_ / _x000A_ (case-insensitive) and newline characters
    cleaned_text = re.sub(r"_x0*00[DA]_", " ", str(text), flags=re.IGNORECASE)
    cleaned_text = cleaned_text.replace("\r", " ").replace("\n", " ")
    cleaned_text = re.sub(r"\s+", " ", cleaned_text).strip()
    
    # Remove any non-ASCII characters to avoid encoding issues
    cleaned_text = ''.join(c for c in cleaned_text if ord(c) < 128)
    
    words = re.split(r"\s+", cleaned_text)
    return cleaned_text if len(words) <= limit else " ".join(words[:limit]) + " ..."

# ───────── Snowflake queries & population helpers ─────────
# populate_boulder_cells(): write GA release counts into the slide table.  
#   - If your pivot logic or GA_COLS change, update this helper accordingly.
def populate_boulder_cells(row, cells, pivot):
    bkey = row["BOULDER_KEY"]
    for idx, rel in enumerate(GA_COLS, start=3):
        val = 0
        if rel != "26.X+":
            # For standard GA releases, get exact value
            val = pivot.at[bkey, rel] if (bkey in pivot.index and rel in pivot.columns) else 0
        else:
            # For 26.X+, sum all columns that start with 26 or higher
            if bkey in pivot.index:
                for col in pivot.columns:
                    # Check if column name starts with a number >= 26
                    if col and str(col).strip() and str(col).split('.')[0].isdigit():
                        if int(str(col).split('.')[0]) >= 26:
                            val += pivot.at[bkey, col]
        cells[idx].text = str(val)

# Actual SQL query for boulder data
SQL_BOULDER = """SELECT boulder_key
     , summary
     , latest_update
     , health_check
     , pillar
     , pillar_sorting
     , boulder_group
FROM CORPANALYTICS_BUSINESS_PROD.SCRATCHPAD_PRDPF.BOULDER_LATEST_DRI_UPDATE_API
WHERE status = 'In Progress'
"""
SQL_FEATURE = """
SELECT strategicinitiative_key AS boulder_key
     , target_ga_release
FROM CORPANALYTICS_BUSINESS_PROD.SCRATCHPAD_PRDPF.FEATURE_API_FULL
WHERE feature_type = 'Type 1'
  AND feature_status <> 'Rejected'
"""

def boulder_by_pillar_load_data() -> tuple[pd.DataFrame, pd.DataFrame]:
    """Load boulder and feature data from Snowflake."""
    logger.info("Connecting to Snowflake...")
    try:
        # Use the dbUtils module to fetch data from Snowflake
        logger.info("Fetching boulder data...")
        boulder_results = execute_snowflake_query(SQL_BOULDER)
        boulders = pd.DataFrame(boulder_results)
        logger.info(f"Retrieved {len(boulders)} boulder records")
        
        # Fetch feature data
        logger.info("Fetching feature data...")
        feature_results = execute_snowflake_query(SQL_FEATURE)
        features = pd.DataFrame(feature_results)
        logger.info(f"Retrieved {len(features)} feature records")
        
        if features.empty:
            logger.warning("No feature data retrieved from Snowflake, falling back to mock data")
            raise ValueError("No feature data retrieved from Snowflake")
            
        return boulders, features
    except Exception as e:
        logger.error(f"Error connecting to Snowflake or retrieving data: {e}")
        logger.info("Falling back to mock data")
        
        # Mock data for Boulder - using pandas DataFrame with the same column structure
        boulder_data = {
            'BOULDER_KEY': ['BDR001', 'BDR002', 'BDR003', 'BDR004', 'BDR005'],
            'SUMMARY': ['Boulder 1 Summary', 'Boulder 2 Summary', 'Boulder 3 Summary', 'Boulder 4 Summary', 'Boulder 5 Summary'],
            'LATEST_UPDATE': ['Latest update for Boulder 1', 'Latest update for Boulder 2', 'Latest update for Boulder 3', 'Latest update for Boulder 4', 'Latest update for Boulder 5'],
            'HEALTH_CHECK': ['Green', 'Yellow', 'Green', 'Red', 'Green'],
            'PILLAR': ['Engineering', 'Engineering', 'Product', 'Product', 'UX'],
            'PILLAR_SORTING': [1, 1, 2, 2, 3],
            'BOULDER_GROUP': ['Group A', 'Group A', 'Group B', 'Group B', 'Group C']
        }
        
        # Mock data for Features - using pandas DataFrame with the same column structure
        feature_data = {
            'BOULDER_KEY': ['BDR001', 'BDR001', 'BDR002', 'BDR003', 'BDR003', 'BDR004', 'BDR005'],
            'TARGET_GA_RELEASE': ['25.03', '25.07', '25.11', '25.03', '26.X+', '25.07', '25.11']
        }
        
        # Print message indicating we're using mock data
        print("WARNING: Using mock data instead of Snowflake connection")
        
        # Create DataFrames from mock data
        boulders_df = pd.DataFrame(boulder_data)
        features_df = pd.DataFrame(feature_data)
        
        # Verify mock data is not empty
        if features_df.empty:
            raise ValueError("Mock feature data is empty!")
            
        return boulders_df, features_df

# ───────── Helpers ─────────
def vmerge_down(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    vMerge = OxmlElement('w:vMerge')
    vMerge.set(qn('w:val'), 'restart')
    tcPr.append(vMerge)

def vmerge_continue(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    vMerge = OxmlElement('w:vMerge')
    tcPr.append(vMerge)

def clone_slide(prs: Presentation, idx: int):
    """Clone *idx* slide from *prs* to preserve the background/layout of template."""
    src = prs.slides[idx]
    # Use the same layout as source slide to preserve background
    layout = src.slide_layout
    dst = prs.slides.add_slide(layout)
    
    # Copy all shapes to maintain formatting and background
    for shp in src.shapes:
        el = deepcopy(shp.element)
        dst.shapes._spTree.insert_element_before(el, 'p:extLst')
    return dst

# ───────── Table population ─────────
# populate_table(): Insert all boulder rows for a pillar into *tbl*.  
#   - Adds/removes rows to exactly match data.  
#   - Merges Boulder_Group cells vertically for contiguous blocks.  
#   - Sets font sizes & colours.  
#   - Customise cell styles (font size/colour, alignment, etc.) in this function.
def populate_table(tbl, df: pd.DataFrame, pivot: pd.DataFrame):
    # capture template row xml for cloning
    template_tr = tbl._tbl.tr_lst[1] if len(tbl.rows) > 1 else None
    # clear body rows (keep header row)
    while len(tbl.rows) > 1:
        tbl._tbl.remove(tbl.rows[1]._tr)

    for _, row in df.iterrows():
        # robust row insertion: use add_row if available, else clone template or add blank row
        if hasattr(tbl, 'add_row'):
            new_row = tbl.add_row()
            cells = new_row.cells
        else:
            if template_tr is not None:
                new_tr = deepcopy(template_tr)
                tbl._tbl.append(new_tr)
            else:
                tbl._tbl.add_tr(0)
            cells = tbl.rows[len(tbl.rows)-1].cells

        # Guard against None values in table cells
        cells[0].text = str(row.get("BOULDER_GROUP", "") or "")
        # Ensure Boulder_Group cell font is white (keep existing background color)
        for para in cells[0].text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)  # ➡️ Change RGB here for different text colour
        cells[1].text = str(row.get("HEALTH_CHECK", "") or "")
        cells[2].text = str(row.get("SUMMARY", "") or "")
        
        # Use the populate_boulder_cells function to fill in the GA release columns
        populate_boulder_cells(row, cells, pivot)

        # Clean and truncate the latest update text
        latest_update = str(row.get("LATEST_UPDATE", "") or "")
        cells[7].text = truncate_words(latest_update, MAX_WORDS)

        # Set font sizes for all cells
        for i, c in enumerate(cells):
            size = 9  # Default size
            if i == 7:  # Latest_Update column
                size = 8
            for para in c.text_frame.paragraphs:
                para.font.size = Pt(size)

    # vertical merge Boulder_Group column
    col, start = 0, 1
    while start < len(tbl.rows):
        end, text = start + 1, tbl.cell(start, col).text
        while end < len(tbl.rows) and tbl.cell(end, col).text == text:
            end += 1
        if end - start > 1:
            vmerge_down(tbl.cell(start, col))
            for r in range(start + 1, end):
                vmerge_continue(tbl.cell(r, col))
                tbl.cell(r, col).text = ""
        start = end

# ───────── Build deck ─────────
# build_deck():  Orchestrates slide creation – one slide per Pillar.
#   • Cleans feature data & builds pivot for GA counts.  
#   • Clones template slide, sets title, removes placeholder header, fills table.
#   • Deletes the original template slide before saving output.
#   - To change slide ordering or add additional elements, edit this function.
def boulder_by_pillar_build_deck(template: str, output: str,
               boulders: pd.DataFrame, features: pd.DataFrame):
    prs = Presentation(template)
    src_slide_idx = 0  # template slide (slide 1)

    # Dynamically extract GA release columns from the template's header row
    try:
        template_tbl = next(
            s for s in prs.slides[src_slide_idx].shapes if s.has_table
        ).table
        header_cells = template_tbl.rows[0].cells
        dynamic_ga_cols: list[str] = []
        for cell in header_cells[3:]:  # skip Boulder_Group, Health_Check, Summary columns
            txt = cell.text.strip()
            # Stop at the Latest_Update column or any blank cell
            if re.match(r"latest\s*_?update", txt, re.IGNORECASE) or not txt:
                break
            dynamic_ga_cols.append(txt)
        global GA_COLS
        GA_COLS = dynamic_ga_cols or GA_COLS  # fallback to default if none found
    except Exception as ex:
        logger.warning(
            "Unable to infer GA_COLS from PPT template; using default list. Error: %s", ex
        )

    # clean GA release (strip '*')
    features["TARGET_GA_RELEASE"] = (
        features["TARGET_GA_RELEASE"].str.replace("*", "", regex=False).str.strip()
    )

    # feature counts pivot
    pivot = (
        features
        .groupby(["BOULDER_KEY", "TARGET_GA_RELEASE"])
        .size()
        .unstack(fill_value=0)
    )
    for col in GA_COLS:
        if col not in pivot.columns:
            pivot[col] = 0

    # sort pillars
    boulders = boulders.sort_values(["PILLAR_SORTING", "PILLAR"], na_position="last")

    slides_to_keep = []
    
    for pillar, df_p in boulders.groupby("PILLAR", sort=False):
        slide = clone_slide(prs, src_slide_idx)
        slides_to_keep.append(slide)
        
        # Set title to actual pillar name
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 1:  # TITLE placeholder
                ph.text = pillar
                break

        # Remove the template header text box (e.g. 'Pillar_name') if it still exists
        for shp in list(slide.shapes):
            if shp.has_text_frame and 'pillar_name' in shp.text_frame.text.lower():
                # Remove the shape from the slide
                sp = shp._element
                sp.getparent().remove(sp)

        try:
            tbl = next(s for s in slide.shapes if s.has_table).table
            populate_table(tbl, df_p.reset_index(drop=True), pivot)
        except StopIteration:
            logger.warning(f"Skipping slide for pillar '{pillar}' because no table was found on the template slide.")
            continue

    # Remove template slide (first slide)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])
    
    prs.save(output)  # Final .pptx written here – change *output* path in main() if needed

# ───────── main ─────────
# Entry point: loads data from Snowflake, builds deck, and handles errors/logging.
#   - To run for a different date range or environment, adjust SQL_BOULDER / SQL_FEATURE.
def main():
    boulders, features = boulder_by_pillar_load_data()
    out_file = f"Boulder Executive Update {datetime.now(timezone.utc):%Y%m%d_%H%M%S}.pptx"
    output   = os.path.join(SCRIPT_DIR, out_file)

    try:
        boulder_by_pillar_build_deck(TEMPLATE, output, boulders, features)
        print(f"✅ Generated {output}")
    except Exception:
        traceback.print_exc()
        raise

if __name__ == "__main__":
    main()
